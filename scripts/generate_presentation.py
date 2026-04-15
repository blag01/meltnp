import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import os

prs = Presentation()

# Helper to add standard slide
def add_slide(title, content, image_paths=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    # Process text content
    if isinstance(content, list):
        for i, point in enumerate(content):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(18)
    
    # Insert multiple images if provided
    if image_paths:
        if isinstance(image_paths, str):
            image_paths = [image_paths]
            
        # Distribute horizontally if multiple
        num_images = len(image_paths)
        start_left = 5.0 if num_images == 1 else 1.0
        width = 4.5 if num_images == 1 else 4.0
        top = 2.5 if num_images == 1 else 4.0
        
        for i, path in enumerate(image_paths):
            if os.path.exists(path):
                left_pos = start_left + (i * width)
                slide.shapes.add_picture(path, Inches(left_pos), Inches(top), width=Inches(width))

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Scaling Transformer Neural Processes"
subtitle.text = "Robustness Benchmarking and Test-Time Adaptation"

# Slide 2: Core Problem
add_slide("The Core Problem", [
    "Neural Processes (NPs) are incredibly powerful at few-shot function regression, but they struggle heavily with Out-Of-Distribution (OOD) data.",
    "Small perturbations (noise, sensor failures, structural bias) often destroy regression coherence.",
    "Our goal: Create a framework capable of measuring robustness under distribution shift."
])

# Slide 3: Model Architecture
add_slide("Model Architecture: Transformer Neural Process", [
    "Moving from flat mean-aggregation (classical NP) to deep Transformer Neural Processes (TNP) (L=3 stacked Transformers with residual connections).",
    "Optional Latent TNP path via Latent Encoders.",
    "Targeted regularized ELBO sampling for stochastic uncertainty bands."
], ["results/tnp/10/gp_robust/weights.png"])

# Slide 4: Structured Benchmarking Suite
add_slide("Structured Benchmarking Suite", [
    "Automated pipeline measuring NLL, MSE, and ECE across 10 corruption levels (0.0 to 2.0).",
    "The x-axis is the raw corruption parameter, not a normalised score:",
    " - Gaussian Noise: sigma (std dev of additive noise, e.g. 2.0 = noise as large as the signal)",
    " - Bias: constant offset added to all y values",
    " - Heteroskedastic: scale factor s of input-dependent noise (noise = s * |x|)",
    " - Outlier: fraction of context replaced with extreme values",
    " - Covariate shift: translation of the test x distribution",
    " - Warp: degree of nonlinear axis distortion",
], ["results/tnp/10/plots/outlier/sinusoid_nll.png"])

# Slide 5: Test-Time Adaptation Concepts
add_slide("Test-Time Adaptation via Pseudo-Likelihood Optimization", [
    "Adapting to corruption at inference time without access to ground-truth labels.",
    "A small MLP is prepended to the frozen TNP and optimized end-to-end using pseudo-likelihood.",
    "The MLP learns a denoising transformation: the TNP's own NLL serves as the only supervision signal."
])

# Slide 6: Visualizing the Adaptation
add_slide("Predictive Quality: Before vs After Adaptation", [
    "Comparing the TNP's predictive distribution before and after the MLP optimization.",
    "Left: High NLL under heavy corruption; wide, uninformative uncertainty bands.",
    "Right: After adaptation, uncertainty concentrates around the inferred function."
], [
    "results/tnp/10/test_time_adaptation/before_mlp.png",
    "results/tnp/10/test_time_adaptation/after_mlp.png"
])

# Slide 7: Adaptation Descent Curve
add_slide("NLL Minimization during Adaptation", [
    "NLL decreases rapidly within the first 50 gradient steps.",
    "The MLP converges to a denoising transformation consistent with the TNP's training prior."
], ["results/tnp/10/test_time_adaptation/optimization_curve_mlp.png"])

# Slide 8: Stochastic Gradient Langevin Dynamics (SGLD)
add_slide("Regularizing Adaptation with SGLD", [
    "Standard gradient descent can converge to degenerate local minima (e.g. constant zero predictions).",
    "SGLD injects calibrated Gaussian noise into parameter updates to encourage exploration of the loss landscape.",
    "Varying noise scale sigma controls the trade-off between exploration and convergence stability."
], ["results/tnp/10/tta_budget/budget_sinusoid_Heteroskedastic_s1.0.png"])

Path("assets").mkdir(exist_ok=True)
output_path = "assets/Neural_Processes_Robustness_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
