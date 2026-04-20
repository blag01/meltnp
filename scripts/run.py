"""
run.py — single entry point for the neural-processes-distribution-shift project.

Usage:
    uv run python scripts/run.py help
    uv run python scripts/run.py sweep [--z-dims none 16] [--clean] [--plots-only] [--no-train] [--no-bench] [--no-extra]
    uv run python scripts/run.py presentation
"""
import subprocess
import sys
import os
import torch
from pathlib import Path

# Ensure src/ is importable
PROJECT_ROOT = Path(__file__).resolve().parents[1]
SRC_DIR = PROJECT_ROOT / "src"
os.environ["PYTHONPATH"] = str(SRC_DIR) + os.pathsep + os.environ.get("PYTHONPATH", "")
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))


# ──────────────────────────────────────────────
# Sweep logic
# ──────────────────────────────────────────────

def run_training_phase(experiments, z_dims):
    """Phase 1: Train all models."""
    for z_dim in z_dims:
        for dataset, robust, num_context in experiments:
            mode = "robust" if robust else "vanilla"
            root = "results/tnp" if z_dim is None else f"results/z{z_dim}tnp"
            output_dir = Path(f"{root}/{num_context}/{dataset}_{mode}")
            output_dir.mkdir(parents=True, exist_ok=True)
            weights_path = output_dir / "weights.pt"

            cmd = [
                sys.executable, str(Path(__file__).parent / "train.py"),
                "--dataset", dataset,
                "--num-context", str(num_context),
                "--epochs", "1000",
                "--output", str(weights_path),
            ]
            if robust:
                cmd.append("--robust")
            if z_dim is not None:
                cmd.extend(["--z-dim", str(z_dim)])

            print(f"\n>>> [Train] {dataset}_{num_context} ({mode})")
            subprocess.run(cmd, check=True)


def run_benchmarking_phase(experiments, z_dims):
    """Phase 2 & 3: Benchmark and generate plots."""
    from np_shift import AttentionNeuralProcess, run_stress_test, plot_robustness_curves

    print("\nStarting benchmarking phase...")
    shift_types = ["noise", "bias", "hetero", "warp", "outlier", "covariate"]
    groups = sorted(
        set((d, c, z) for d, _, c in experiments for z in z_dims),
        key=lambda x: (x[0], x[1], x[2] if x[2] is not None else -1),
    )
    all_results = {g: {st: {} for st in shift_types} for g in groups}

    for z_dim in z_dims:
        for dataset, robust, num_context in experiments:
            mode = "robust" if robust else "vanilla"
            root = "results/tnp" if z_dim is None else f"results/z{z_dim}tnp"
            model_name = f"{dataset}_{mode}_{num_context}"
            weights_path = Path(f"{root}/{num_context}/{dataset}_{mode}/weights.pt")

            if not weights_path.exists():
                continue

            model = AttentionNeuralProcess(z_dim=z_dim)
            model.load_state_dict(torch.load(weights_path, weights_only=True))
            model.eval()

            print(f"Benchmarking: {model_name}...")
            for st in shift_types:
                all_results[(dataset, num_context, z_dim)][st][model_name] = run_stress_test(
                    model, dataset, st, num_context=num_context)

            if not robust:
                for tta_method in ["mlp", "reweight", "latent",
                                   "mlp_sgld_0.01", "mlp_sgld_0.05",
                                   "mlp_sgld_0.1", "mlp_sgld_0.2"]:
                    tta_name = f"{model_name}_tta_{tta_method}"
                    print(f"  TTA: {tta_name}...")
                    for st in shift_types:
                        all_results[(dataset, num_context, z_dim)][st][tta_name] = run_stress_test(
                            model, dataset, st, adapt_method=tta_method, num_context=num_context)

    print("Generating robustness curves...")
    for ds, ctx, z_dim in groups:
        root = "results/tnp" if z_dim is None else f"results/z{z_dim}tnp"
        plot_dir = Path(f"{root}/{ctx}/plots")
        for st in shift_types:
            if all_results[(ds, ctx, z_dim)][st]:
                st_dir = plot_dir / st
                st_dir.mkdir(parents=True, exist_ok=True)
                plot_robustness_curves(all_results[(ds, ctx, z_dim)][st], str(st_dir), file_prefix=ds)
    print("All robustness plots saved.")


# ──────────────────────────────────────────────
# Subcommand handlers
# ──────────────────────────────────────────────

HELP_TEXT = """\
Neural Processes under Distribution Shift — run.py

Shorthand commands (recommended):
  train        Full run from scratch: train → benchmark → TTA → budget curves.
  plot         Re-run benchmark + extras on existing weights (skip training).
  present      Generate the PowerPoint slide deck.

Full command:
  sweep        Run experiment pipeline with custom flags (see below).
  help         Show this message.

Sweep flags:
  --z-dims none 16   Model variants (none = deterministic TNP, 16 = Latent TNP).
  --clean            Wipe results/ before starting.
  --plots-only       Skip training; re-run benchmark + extras on existing weights.
  --no-train         Skip training phase only.
  --no-bench         Skip benchmarking phase only.
  --no-extra         Skip TTA visual and budget scripts.

Examples:
  uv run python scripts/run.py train
  uv run python scripts/run.py plot
  uv run python scripts/run.py sweep --no-extra --z-dims none
"""


def cmd_sweep(argv):
    import argparse
    parser = argparse.ArgumentParser(prog="run.py sweep")
    parser.add_argument("--no-train",   action="store_true")
    parser.add_argument("--no-bench",   action="store_true")
    parser.add_argument("--no-extra",   action="store_true")
    parser.add_argument("--plots-only", action="store_true")
    parser.add_argument("--z-dims", nargs="+", default=["none", "16"])
    parser.add_argument("--clean",      action="store_true")
    args = parser.parse_args(argv)

    if args.plots_only:
        args.no_train = True

    if args.clean:
        import shutil
        if Path("results").exists():
            print(">>> [Cleanup] Wiping results/ ...")
            shutil.rmtree("results")

    z_dims = [None if z.lower() == "none" else int(z) for z in args.z_dims]

    datasets      = ["gp", "sinusoid"]
    robust_flags  = [False, True]
    context_sizes = [10, 20, 40]
    experiments   = [(ds, r, ctx)
                     for ctx in context_sizes
                     for ds  in datasets
                     for r   in robust_flags]

    if not args.no_train:
        run_training_phase(experiments, z_dims)
    if not args.no_bench:
        run_benchmarking_phase(experiments, z_dims)
    if not args.no_extra:
        scripts = Path(__file__).parent
        for z in z_dims:
            extra = ["--z-dim", str(z)] if z is not None else []
            print(f"\n>>> [Extra] TTA Prototypes (z_dim={z})")
            subprocess.run([sys.executable, str(scripts / "test_time_adapt.py")] + extra, check=True)
            print(f"\n>>> [Extra] TTA Budget Curves (z_dim={z})")
            subprocess.run([sys.executable, str(scripts / "tta_budget.py")] + extra, check=True)

    print("\nSweep complete.")


def cmd_presentation(_argv):
    import collections, collections.abc, os
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("Error: python-pptx not installed. Run:  uv sync")
        sys.exit(1)

    prs = Presentation()

    def add_slide(title, content, image_paths=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        for i, point in enumerate(content):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(18)
        if image_paths:
            if isinstance(image_paths, str):
                image_paths = [image_paths]
            n = len(image_paths)
            start_left, width, top = (5.0, 4.5, 2.5) if n == 1 else (1.0, 4.0, 4.0)
            for i, path in enumerate(image_paths):
                if os.path.exists(path):
                    slide.shapes.add_picture(path, Inches(start_left + i * width), Inches(top), width=Inches(width))

    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Neural Processes under Distribution Shift"
    s.placeholders[1].text = "Robustness Benchmarking and Test-Time Adaptation"

    add_slide("The Core Problem", [
        "Neural Processes (NPs) are incredibly powerful at few-shot function regression, but they struggle heavily with Out-Of-Distribution (OOD) data.",
        "Small perturbations (noise, sensor failures, structural bias) often destroy regression coherence.",
        "Our goal: Create a framework capable of measuring robustness under distribution shift.",
    ])
    add_slide("Model Architecture: Transformer Neural Process", [
        "Moving from flat mean-aggregation (classical NP) to deep Transformer Neural Processes (TNP) (L=3 stacked Transformers with residual connections).",
        "Optional Latent TNP path via Latent Encoders.",
        "Targeted regularized ELBO sampling for stochastic uncertainty bands.",
    ], ["results/tnp/10/gp_robust/weights.png"])
    add_slide("Structured Benchmarking Suite", [
        "Automated pipeline measuring NLL, MSE, and ECE across 10 corruption levels (0.0 to 2.0).",
        "The x-axis is the raw corruption parameter, not a normalised score:",
        " - Gaussian Noise: sigma (e.g. 2.0 = noise as large as the signal)",
        " - Bias: constant offset added to all y values",
        " - Heteroskedastic: scale factor s (noise = s * |x|)",
        " - Outlier: fraction of context replaced with extreme values",
        " - Covariate shift: translation of the test x distribution",
        " - Warp: degree of nonlinear axis distortion",
    ], ["results/tnp/10/plots/outlier/sinusoid_nll.png"])
    add_slide("Test-Time Adaptation via Pseudo-Likelihood Optimization", [
        "Adapting to corruption at inference time without access to ground-truth labels.",
        "A small MLP is prepended to the frozen TNP and optimized end-to-end using pseudo-likelihood.",
        "The MLP learns a denoising transformation: the TNP's own NLL serves as the only supervision signal.",
    ])
    add_slide("Predictive Quality: Before vs After Adaptation", [
        "Comparing the TNP's predictive distribution before and after the MLP optimization.",
        "Left: High NLL under heavy corruption; wide, uninformative uncertainty bands.",
        "Right: After adaptation, uncertainty concentrates around the inferred function.",
    ], [
        "results/tnp/10/test_time_adaptation/before_mlp.png",
        "results/tnp/10/test_time_adaptation/after_mlp.png",
    ])
    add_slide("NLL Minimization during Adaptation", [
        "NLL decreases rapidly within the first 50 gradient steps.",
        "The MLP converges to a denoising transformation consistent with the TNP's training prior.",
    ], ["results/tnp/10/test_time_adaptation/optimization_curve_mlp.png"])
    add_slide("Regularizing Adaptation with SGLD", [
        "Standard gradient descent can converge to degenerate local minima (e.g. constant zero predictions).",
        "SGLD injects calibrated Gaussian noise into parameter updates to encourage exploration of the loss landscape.",
        "Varying noise scale sigma controls the trade-off between exploration and convergence stability.",
    ], ["results/tnp/10/tta_budget/budget_sinusoid_Heteroskedastic_s1.0.png"])

    Path("assets").mkdir(exist_ok=True)
    out = "assets/Neural_Processes_Robustness_Presentation.pptx"
    prs.save(out)
    print(f"Presentation saved to {out}")


def cmd_help(_argv):
    print(HELP_TEXT)


def cmd_train(_argv):
    """Shorthand: full run from scratch."""
    cmd_sweep(["--clean"])


def cmd_plot(_argv):
    """Shorthand: re-run benchmark + extras on existing weights."""
    cmd_sweep(["--plots-only"])


def cmd_present(argv):
    """Shorthand alias for presentation."""
    cmd_presentation(argv)


COMMANDS = {
    # Shorthands
    "train":        cmd_train,
    "plot":         cmd_plot,
    "present":      cmd_present,
    # Full control
    "sweep":        cmd_sweep,
    "presentation": cmd_presentation,
    "help":         cmd_help,
}


def main():
    if len(sys.argv) < 2 or sys.argv[1] not in COMMANDS:
        print(HELP_TEXT)
        sys.exit(0 if len(sys.argv) < 2 else 1)
    COMMANDS[sys.argv[1]](sys.argv[2:])


if __name__ == "__main__":
    main()
